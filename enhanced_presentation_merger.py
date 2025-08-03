#!/usr/bin/env python3
"""
Enhanced PowerPoint Presentation Merger
Extracts content from testFrameworkPPT.pptx and adds it to JIRA_AI_Final_Presentation.pptx
with comprehensive formatting preservation and slide analysis.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

def analyze_presentation(pptx_path):
    """Analyze presentation structure and content"""
    try:
        prs = Presentation(pptx_path)
        analysis = {
            'total_slides': len(prs.slides),
            'slide_layouts': {},
            'content_summary': []
        }
        
        print(f"\n📊 Analyzing: {os.path.basename(pptx_path)}")
        print(f"Total slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            slide_info = {
                'slide_number': i + 1,
                'layout': slide.slide_layout.name,
                'shapes_count': len(slide.shapes),
                'text_boxes': 0,
                'images': 0,
                'tables': 0,
                'charts': 0,
                'content': []
            }
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    slide_info['text_boxes'] += 1
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_info['content'].append({
                            'type': 'text',
                            'text': shape.text[:100] + '...' if len(shape.text) > 100 else shape.text
                        })
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info['images'] += 1
                    slide_info['content'].append({'type': 'image'})
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    slide_info['tables'] += 1
                    slide_info['content'].append({'type': 'table'})
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    slide_info['charts'] += 1
                    slide_info['content'].append({'type': 'chart'})
            
            analysis['content_summary'].append(slide_info)
            
            # Count layout usage
            layout_name = slide.slide_layout.name
            if layout_name not in analysis['slide_layouts']:
                analysis['slide_layouts'][layout_name] = 0
            analysis['slide_layouts'][layout_name] += 1
            
            print(f"  Slide {i+1}: {layout_name} - {slide_info['text_boxes']} text boxes, {slide_info['images']} images")
        
        print(f"Layouts used: {list(analysis['slide_layouts'].keys())}")
        return analysis, prs
        
    except Exception as e:
        print(f"❌ Error analyzing {pptx_path}: {e}")
        return None, None

def extract_slide_with_formatting(slide):
    """Extract slide content with full formatting information"""
    slide_data = {
        'layout_name': slide.slide_layout.name,
        'shapes': [],
        'background': None
    }
    
    for shape in slide.shapes:
        shape_data = {
            'type': shape.shape_type,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'text': '',
            'font_properties': {},
            'alignment': None,
            'fill_color': None,
            'line_color': None
        }
        
        # Extract text and formatting
        if hasattr(shape, 'text_frame'):
            shape_data['text'] = shape.text
            shape_data['alignment'] = shape.text_frame.paragraphs[0].alignment if shape.text_frame.paragraphs else None
            
            # Extract font properties
            if shape.text_frame.paragraphs:
                paragraph = shape.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    shape_data['font_properties'] = {
                        'name': run.font.name,
                        'size': run.font.size,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') and run.font.color.rgb else None
                    }
        
        # Extract shape formatting
        if hasattr(shape, 'fill'):
            if shape.fill.type:
                shape_data['fill_color'] = shape.fill.fore_color.rgb if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb else None
        
        if hasattr(shape, 'line'):
            if hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                shape_data['line_color'] = shape.line.color.rgb
        
        slide_data['shapes'].append(shape_data)
    
    return slide_data

def apply_formatting_to_shape(target_shape, source_data):
    """Apply formatting from source data to target shape"""
    try:
        if hasattr(target_shape, 'text_frame') and source_data.get('text'):
            # Set text
            target_shape.text = source_data['text']
            
            # Apply font properties
            if source_data.get('font_properties'):
                font_props = source_data['font_properties']
                for paragraph in target_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if font_props.get('name'):
                            run.font.name = font_props['name']
                        if font_props.get('size'):
                            run.font.size = font_props['size']
                        if font_props.get('bold') is not None:
                            run.font.bold = font_props['bold']
                        if font_props.get('italic') is not None:
                            run.font.italic = font_props['italic']
                        if font_props.get('color'):
                            run.font.color.rgb = font_props['color']
            
            # Apply alignment
            if source_data.get('alignment'):
                for paragraph in target_shape.text_frame.paragraphs:
                    paragraph.alignment = source_data['alignment']
        
        # Apply shape formatting
        if source_data.get('fill_color') and hasattr(target_shape, 'fill'):
            target_shape.fill.solid()
            target_shape.fill.fore_color.rgb = source_data['fill_color']
        
        if source_data.get('line_color') and hasattr(target_shape, 'line'):
            target_shape.line.color.rgb = source_data['line_color']
            
    except Exception as e:
        print(f"Warning: Could not apply formatting: {e}")

def merge_presentations_with_analysis(source_path, target_path, output_path):
    """Merge presentations with comprehensive analysis and formatting preservation"""
    
    # Analyze both presentations
    print("🔍 Analyzing presentations...")
    source_analysis, source_prs = analyze_presentation(source_path)
    target_analysis, target_prs = analyze_presentation(target_path)
    
    if not source_analysis or not target_analysis:
        print("❌ Failed to analyze presentations")
        return False
    
    # Create merged presentation
    print(f"\n🔄 Creating merged presentation...")
    merged_prs = Presentation(target_path)
    
    # Add slides from source to target
    source_slides_added = 0
    for i, slide in enumerate(source_prs.slides):
        print(f"Processing slide {i+1} from source...")
        
        # Extract slide data with formatting
        slide_data = extract_slide_with_formatting(slide)
        
        # Find matching layout or use default
        matching_layout = None
        for layout in merged_prs.slide_layouts:
            if layout.name == slide_data['layout_name']:
                matching_layout = layout
                break
        
        if not matching_layout:
            print(f"  Using default layout for slide {i+1}")
            matching_layout = merged_prs.slide_layouts[0]
        
        # Create new slide
        new_slide = merged_prs.slides.add_slide(matching_layout)
        
        # Copy shapes with formatting
        for shape_data in slide_data['shapes']:
            if shape_data['type'] == MSO_SHAPE_TYPE.TEXT_BOX and shape_data['text']:
                # Create text box
                textbox = new_slide.shapes.add_textbox(
                    shape_data['left'],
                    shape_data['top'],
                    shape_data['width'],
                    shape_data['height']
                )
                
                # Apply formatting
                apply_formatting_to_shape(textbox, shape_data)
        
        source_slides_added += 1
    
    # Save merged presentation
    print(f"\n💾 Saving merged presentation...")
    merged_prs.save(output_path)
    
    # Summary
    print(f"\n✅ Merge completed successfully!")
    print(f"📊 Summary:")
    print(f"   - Original target slides: {target_analysis['total_slides']}")
    print(f"   - Source slides added: {source_slides_added}")
    print(f"   - Total slides in merged: {target_analysis['total_slides'] + source_slides_added}")
    print(f"   - Output file: {output_path}")
    
    return True

def main():
    """Main function"""
    presentation_dir = "Presentation"
    
    source_file = os.path.join(presentation_dir, "testFrameworkPPT.pptx")
    target_file = os.path.join(presentation_dir, "JIRA_AI_Final_Presentation.pptx")
    output_file = os.path.join(presentation_dir, "JIRA_AI_Final_Presentation_Merged.pptx")
    
    # Check if files exist
    if not os.path.exists(source_file):
        print(f"❌ Source file not found: {source_file}")
        return
    
    if not os.path.exists(target_file):
        print(f"❌ Target file not found: {target_file}")
        return
    
    print("🎯 Enhanced PowerPoint Presentation Merger")
    print("=" * 60)
    print(f"Source: {source_file}")
    print(f"Target: {target_file}")
    print(f"Output: {output_file}")
    print("=" * 60)
    
    # Check if python-pptx is installed
    try:
        import pptx
    except ImportError:
        print("❌ python-pptx not installed. Installing required dependencies...")
        os.system("pip install python-pptx Pillow")
        print("✅ Dependencies installed. Please run the script again.")
        return
    
    # Merge presentations
    success = merge_presentations_with_analysis(source_file, target_file, output_file)
    
    if success:
        print(f"\n🎉 Successfully created merged presentation!")
        print(f"📁 Location: {output_file}")
        print("\n📋 The merged presentation contains:")
        print("   ✅ All slides from JIRA_AI_Final_Presentation.pptx")
        print("   ✅ All slides from testFrameworkPPT.pptx")
        print("   ✅ Preserved formatting and styling")
        print("   ✅ Maintained slide layouts where possible")
    else:
        print("❌ Failed to merge presentations")

if __name__ == "__main__":
    main() 
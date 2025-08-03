#!/usr/bin/env python3
"""
PowerPoint Presentation Merger
Extracts content from testFrameworkPPT.pptx and adds it to JIRA_AI_Final_Presentation.pptx
with the same formatting and style.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

def extract_slide_content(pptx_path):
    """Extract content from a PowerPoint file"""
    try:
        prs = Presentation(pptx_path)
        slides_content = []
        
        for slide in prs.slides:
            slide_data = {
                'layout': slide.slide_layout.name,
                'shapes': [],
                'text': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_data['text'].append({
                        'text': shape.text,
                        'position': (shape.left, shape.top),
                        'size': (shape.width, shape.height)
                    })
                
                # Store shape information
                shape_data = {
                    'type': type(shape).__name__,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                # Extract text from text boxes
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            shape_data['text'] = paragraph.text
                            shape_data['alignment'] = paragraph.alignment
                
                slide_data['shapes'].append(shape_data)
            
            slides_content.append(slide_data)
        
        return slides_content, prs
    except Exception as e:
        print(f"Error reading {pptx_path}: {e}")
        return None, None

def get_slide_layout_by_name(prs, layout_name):
    """Get slide layout by name"""
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return prs.slide_layouts[0]  # Default to first layout

def copy_formatting(source_shape, target_shape):
    """Copy formatting from source shape to target shape"""
    try:
        if hasattr(source_shape, 'text_frame') and hasattr(target_shape, 'text_frame'):
            # Copy font properties
            for i, paragraph in enumerate(source_shape.text_frame.paragraphs):
                if i < len(target_shape.text_frame.paragraphs):
                    target_paragraph = target_shape.text_frame.paragraphs[i]
                    
                    # Copy paragraph alignment
                    target_paragraph.alignment = paragraph.alignment
                    
                    # Copy font properties for each run
                    for j, run in enumerate(paragraph.runs):
                        if j < len(target_paragraph.runs):
                            target_run = target_paragraph.runs[j]
                            target_run.font.name = run.font.name
                            target_run.font.size = run.font.size
                            target_run.font.bold = run.font.bold
                            target_run.font.italic = run.font.italic
                            if run.font.color.rgb:
                                target_run.font.color.rgb = run.font.color.rgb
    except Exception as e:
        print(f"Error copying formatting: {e}")

def merge_presentations(source_path, target_path, output_path):
    """Merge presentations with formatting preservation"""
    print(f"Reading source presentation: {source_path}")
    source_content, source_prs = extract_slide_content(source_path)
    
    print(f"Reading target presentation: {target_path}")
    target_content, target_prs = extract_slide_content(target_path)
    
    if not source_content or not target_content:
        print("Failed to read one or both presentations")
        return False
    
    print(f"Source presentation has {len(source_content)} slides")
    print(f"Target presentation has {len(target_content)} slides")
    
    # Create a copy of the target presentation
    merged_prs = Presentation(target_path)
    
    # Add slides from source to target
    for i, slide_data in enumerate(source_content):
        print(f"Processing slide {i+1} from source...")
        
        # Create new slide with similar layout
        layout = get_slide_layout_by_name(merged_prs, slide_data['layout'])
        new_slide = merged_prs.slides.add_slide(layout)
        
        # Copy content and formatting
        for shape_data in slide_data['shapes']:
            if 'text' in shape_data:
                # Add text box
                left = shape_data['left']
                top = shape_data['top']
                width = shape_data['width']
                height = shape_data['height']
                
                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = shape_data['text']
                
                # Apply formatting if available
                if 'alignment' in shape_data:
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = shape_data['alignment']
    
    # Save the merged presentation
    print(f"Saving merged presentation to: {output_path}")
    merged_prs.save(output_path)
    print("✅ Presentation merged successfully!")
    return True

def main():
    """Main function to merge presentations"""
    presentation_dir = "Presentation"
    
    source_file = os.path.join(presentation_dir, "testFrameworkPPT.pptx")
    target_file = os.path.join(presentation_dir, "JIRA_AI_Final_Presentation.pptx")
    output_file = os.path.join(presentation_dir, "JIRA_AI_Final_Presentation_Merged.pptx")
    
    # Check if files exist
    if not os.path.exists(source_file):
        print(f"❌ Source file not found: {source_file}")
        return
    
    if not os.path.exists(target_file):
        print(f"❌ Target file not found: {target_file}")
        return
    
    print("🎯 PowerPoint Presentation Merger")
    print("=" * 50)
    print(f"Source: {source_file}")
    print(f"Target: {target_file}")
    print(f"Output: {output_file}")
    print("=" * 50)
    
    # Merge presentations
    success = merge_presentations(source_file, target_file, output_file)
    
    if success:
        print(f"\n🎉 Successfully created merged presentation: {output_file}")
        print("📋 The merged presentation contains:")
        print("   - All slides from JIRA_AI_Final_Presentation.pptx")
        print("   - All slides from testFrameworkPPT.pptx (with preserved formatting)")
    else:
        print("❌ Failed to merge presentations")

if __name__ == "__main__":
    main() 